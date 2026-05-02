#!/usr/bin/env python3
"""Build the 3-slide Lemma deck for the Beta Super Hackathon submission.

Style: Apple keynote / A24 — black background, cream Fraunces serif title, Inter
sans for body. Mirrors the Lemma app's visual identity so deck + product + video
read as one brand.
"""
import os
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Paths
ROOT = Path("/Users/touyuumiyabi/Desktop/lemma ai")
VIDEO = ROOT / "public" / "demo-videos" / "final-2min-compressed.mp4"
OUT = ROOT / "lemma-deck.pptx"

# Brand colors
INK = RGBColor(0x0E, 0x0E, 0x0F)
CREAM = RGBColor(0xF4, 0xEF, 0xE6)
BONE = RGBColor(0xC9, 0xC1, 0xB1)
MUTED = RGBColor(0x8A, 0x82, 0x75)
EMBER = RGBColor(0xE8, 0x5D, 0x24)

# Fonts (Mac defaults; PPT will substitute if missing on Windows)
SERIF = "Georgia"   # closest mac default to Fraunces
SANS = "Helvetica"

# Slide canvas — 16:9
WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)


def set_bg_black(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = INK


def add_text(slide, x, y, w, h, text, *, font=SANS, size=18, color=CREAM,
             bold=False, italic=False, align=PP_ALIGN.LEFT, tracking=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_eyebrow(slide, x, y, text):
    return add_text(
        slide, x, y, Inches(8), Inches(0.4),
        text.upper(), font=SANS, size=11, color=MUTED, bold=True,
    )


def add_divider(slide, x, y, w):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(6350))
    line.fill.solid()
    line.fill.fore_color.rgb = BONE
    line.line.fill.background()
    return line


def slide1_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_black(s)

    add_eyebrow(s, Inches(0.7), Inches(0.6), "Lemma · Beta Super Hackathon · May 2 2026")

    # Big serif headline
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "We are "
    r1.font.name = SERIF
    r1.font.size = Pt(72)
    r1.font.color.rgb = CREAM
    r2 = p.add_run()
    r2.text = "Lemma."
    r2.font.name = SERIF
    r2.font.size = Pt(72)
    r2.font.italic = True
    r2.font.color.rgb = CREAM

    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r3 = p2.add_run()
    r3.text = "Adaptive listing videos for short-term rental hosts."
    r3.font.name = SERIF
    r3.font.size = Pt(28)
    r3.font.italic = True
    r3.font.color.rgb = BONE

    # Two member cards
    add_divider(s, Inches(0.7), Inches(4.4), Inches(12))

    # Xiya
    add_eyebrow(s, Inches(0.7), Inches(4.7), "Founder · Product + Engineering")
    add_text(s, Inches(0.7), Inches(5.05), Inches(6), Inches(0.7),
             "Xiya Tang", font=SERIF, size=36, color=CREAM)
    add_text(s, Inches(0.7), Inches(5.85), Inches(6), Inches(1.5),
             "Designer-engineer. Built and shipped Lemma's full stack — Director Agent on Claude, "
             "Seedance 2.0 video pipeline, Butterbase backend, embeddable widget.",
             font=SANS, size=14, color=BONE)

    # Member
    add_eyebrow(s, Inches(7), Inches(4.7), "Co-founder · Business + GTM")
    add_text(s, Inches(7), Inches(5.05), Inches(6), Inches(0.7),
             "Shourya Premkumar", font=SERIF, size=36, color=CREAM)
    add_text(s, Inches(7), Inches(5.6), Inches(6), Inches(0.4),
             "American High School", font=SERIF, size=16, italic=True, color=BONE)
    add_text(s, Inches(7), Inches(6.05), Inches(6), Inches(1.3),
             "Founder, STEAD Marketing — AI-driven online marketing optimization for brands. "
             "Direct experience driving brand conversions through targeted content strategy.",
             font=SANS, size=13, color=BONE)

    # Bottom tag
    add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
             "Why us · we are the user. We've hosted, built listing pages, "
             "and watched conversion die on static photo galleries.",
             font=SANS, size=10, color=MUTED, italic=True)


def slide2_product(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_black(s)

    add_eyebrow(s, Inches(0.7), Inches(0.6), "Product Overview")

    # The one-liner
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "One stay, "
    r1.font.name = SERIF
    r1.font.size = Pt(64)
    r1.font.color.rgb = CREAM
    r2 = p.add_run()
    r2.text = "told four ways."
    r2.font.name = SERIF
    r2.font.size = Pt(64)
    r2.font.italic = True
    r2.font.color.rgb = CREAM

    add_text(s, Inches(0.7), Inches(3.05), Inches(11), Inches(0.6),
             "Property managers upload photos once. Our agent generates four cinematic "
             "films — one per traveler — and embeds them on any direct-booking page.",
             font=SERIF, size=18, italic=True, color=BONE)

    # Two columns: Problem / Solution
    PROB_X = Inches(0.7)
    SOL_X = Inches(7.0)
    BOX_TOP = Inches(4.2)

    # Problem
    add_eyebrow(s, PROB_X, BOX_TOP, "The problem")
    add_text(s, PROB_X, Inches(4.55), Inches(5.8), Inches(0.5),
             "One generic walkthrough.", font=SERIF, size=22, color=CREAM)
    add_text(s, PROB_X, Inches(5.05), Inches(5.8), Inches(2),
             "Every guest — couple, family, remote worker, business — sees the same listing video. "
             "Direct-booking pages convert at 4×worse than Airbnb because they all feel generic. "
             "Hosts have no tool to fix this without a videographer.",
             font=SANS, size=13, color=BONE)

    # Solution
    add_eyebrow(s, SOL_X, BOX_TOP, "What we built")
    add_text(s, SOL_X, Inches(4.55), Inches(5.8), Inches(0.5),
             "Adaptive video, in 90 seconds.", font=SERIF, size=22, color=CREAM)
    sol_bullets = [
        "Director Agent on Claude reads the home and proposes bespoke personas",
        "Seedance 2.0 renders 4 cinematic films, one per persona",
        "One <script> tag embeds an adaptive widget on the host's site",
        "Live analytics on Butterbase: who's actually watching what",
    ]
    y = Inches(5.05)
    for b in sol_bullets:
        add_text(s, SOL_X, y, Inches(5.8), Inches(0.4), "·", font=SANS, size=14, color=EMBER)
        add_text(s, SOL_X + Inches(0.25), y, Inches(5.55), Inches(0.4),
                 b, font=SANS, size=12.5, color=BONE)
        y += Inches(0.4)

    # Stack tag
    add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
             "Stack · Claude (Director Agent) + Seedance 2.0 (video) + Butterbase (Postgres + storage + REST)",
             font=SANS, size=10, color=MUTED, italic=True)


def slide3_demo(prs, video_path: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_black(s)

    add_eyebrow(s, Inches(0.7), Inches(0.45), "Demo · 2 minutes")
    add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.6),
             "See it work.", font=SERIF, size=36, color=CREAM)

    # Embed the video. Generate a poster frame from the source so PPT shows
    # something instead of black before play.
    poster = video_path.with_suffix(".jpg")
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-loglevel", "error", "-ss", "0:08", "-i", str(video_path),
             "-frames:v", "1", "-q:v", "2", str(poster)],
            check=True,
        )
    except Exception:
        poster = None

    # Place video centered, large
    vid_w = Inches(11.5)
    vid_h = Inches(11.5 * 9 / 16)  # ~6.47"
    vid_x = (WIDTH - vid_w) / 2
    vid_y = Inches(1.7)

    s.shapes.add_movie(
        str(video_path),
        vid_x, vid_y, vid_w, vid_h,
        poster_frame_image=str(poster) if poster else None,
        mime_type="video/mp4",
    )

    add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
             "lemma.app  ·  Built with Claude · Seedance 2.0 · Butterbase  ·  #betahacks",
             font=SANS, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def main():
    if not VIDEO.exists():
        raise SystemExit(f"Video not found: {VIDEO}")

    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT

    slide1_team(prs)
    slide2_product(prs)
    slide3_demo(prs, VIDEO)

    prs.save(OUT)
    print(f"✓ Saved: {OUT}")
    sz = os.path.getsize(OUT) / (1024 * 1024)
    print(f"  size: {sz:.1f} MB")


if __name__ == "__main__":
    main()
